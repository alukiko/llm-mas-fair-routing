import os
from dataclasses import dataclass
from pathlib import Path
from typing import Annotated, List, Optional

from fastmcp import Context, FastMCP
from markitdown import MarkItDown
from pptx import Presentation
from pydantic import Field

from automas.mcp.servers.content_utils import create_temp_directory, hash_file, truncate_text
from automas.mcp.servers.document.image_utils import (
    format_image_section,
    get_cached_images,
    save_cached_images,
)

CACHE_DIR = Path.home() / ".automas" / "pptx_cache"
CACHE_DIR.mkdir(parents=True, exist_ok=True)

md = MarkItDown()

pptx_server = FastMCP("pptx-reader")


@dataclass
class ExtractedPPTXImage:
    path: str
    slide_number: int
    shape_name: str
    width: float
    height: float


async def extract_images_from_pptx(pptx_path: str, ctx: Context) -> List[ExtractedPPTXImage]:
    pptx_hash = hash_file(pptx_path)

    cached = await get_cached_images(CACHE_DIR, pptx_hash, ExtractedPPTXImage)
    if cached is not None:
        await ctx.info(f"Using cached images ({len(cached)} images)")
        return cached

    output_dir = create_temp_directory("pptx_images")

    images = []

    try:
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)

        await ctx.info(f"Extracting images from {total_slides} slides")

        for slide_idx, slide in enumerate(prs.slides, start=1):
            if total_slides > 1:
                await ctx.report_progress(progress=slide_idx, total=total_slides)
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        ext = image.ext
                        if not ext.startswith("."):
                            ext = f".{ext}"

                        filename = f"slide{slide_idx}_shape{shape_idx}{ext}"
                        image_path = output_dir / filename

                        with open(image_path, "wb") as f:
                            f.write(image.blob)

                        images.append(
                            ExtractedPPTXImage(
                                path=str(image_path),
                                slide_number=slide_idx,
                                shape_name=shape.name,
                                width=shape.width,
                                height=shape.height,
                            )
                        )
                    except Exception:
                        pass

    except Exception as e:
        raise Exception(f"Error extracting images from PPTX: {e}")

    await save_cached_images(CACHE_DIR, pptx_hash, images)

    if images:
        await ctx.info(f"Extracted {len(images)} images from PPTX")

    return images


@pptx_server.tool
async def read_pptx(
    file_path: Annotated[str, Field(description="Path to the PPTX file")],
    ctx: Context,
    max_lines: Annotated[
        Optional[int],
        Field(description="Maximum number of text lines to return. If None, returns all lines"),
    ] = 1000,
) -> str:
    """
    Read a PPTX file and return its text content.

    Args:
        file_path: Path to the PPTX file (PowerPoint presentation)
        max_lines: Maximum number of text lines to return (default: 1000)

    Returns:
        Text content of the PPTX file
    """
    try:
        expanded_path = os.path.expanduser(file_path)
        file_name = os.path.basename(file_path)

        await ctx.info(f"Reading PPTX: {file_name}")

        text_content = md.convert(expanded_path).text_content
        text_content = truncate_text(text_content, max_lines)

        images = await extract_images_from_pptx(expanded_path, ctx)

        text_content += format_image_section(
            images, [("path", "Path"), ("slide_number", "Slide"), ("shape_name", "Shape")]
        )

        return text_content

    except Exception as e:
        await ctx.error(f"Failed to read PPTX: {e}")
        return f"Error reading PPTX: {e}"
